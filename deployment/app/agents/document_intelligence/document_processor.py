"""Document processing with multi-format support and OCR.

This module handles loading, processing, and chunking of various
document formats including PDF, TXT, DOCX, PPTX, and images.

Following Enterprise Development Standards:
- Software Architect: Modular loader design
- Security Architect: Safe temporary file handling
- Data Architect: Consistent chunk metadata
- Software Engineer: Type-safe with error handling
"""

import io
import logging
import os
import tempfile
from pathlib import Path
from typing import Any

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Handles document loading, OCR, and chunking for multiple formats.

    Supported formats:
    - PDF (.pdf) - via PyPDFLoader
    - Text (.txt) - via TextLoader
    - Word (.docx, .doc) - via Docx2txtLoader
    - PowerPoint (.pptx, .ppt) - via python-pptx
    - Images (.png, .jpg, .jpeg) - via pytesseract OCR
    """

    SUPPORTED_EXTENSIONS = {
        ".pdf", ".txt", ".docx", ".doc", ".pptx", ".ppt",
        ".png", ".jpg", ".jpeg"
    }

    def __init__(
        self,
        chunk_size: int | None = None,
        chunk_overlap: int | None = None,
    ) -> None:
        """Initialize the document processor.

        Args:
            chunk_size: Size of text chunks (default from env or 1000)
            chunk_overlap: Overlap between chunks (default from env or 200)
        """
        self.chunk_size = chunk_size or int(os.getenv("DOC_CHUNK_SIZE", "1000"))
        self.chunk_overlap = chunk_overlap or int(os.getenv("DOC_CHUNK_OVERLAP", "200"))

        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", " ", ""],
        )

        self._setup_tesseract()

    def _setup_tesseract(self) -> None:
        """Configure Tesseract OCR path for Windows."""
        tesseract_cmd = os.getenv("TESSERACT_CMD")
        if tesseract_cmd:
            try:
                import pytesseract
                pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
                logger.info(f"Tesseract path set to: {tesseract_cmd}")
            except ImportError:
                logger.warning("pytesseract not installed, OCR will not be available")

    def is_supported(self, filename: str) -> bool:
        """Check if file type is supported.

        Args:
            filename: Filename to check

        Returns:
            True if file type is supported
        """
        ext = Path(filename).suffix.lower()
        return ext in self.SUPPORTED_EXTENSIONS

    def get_file_type(self, filename: str) -> str:
        """Get the file type category for a filename.

        Args:
            filename: Filename to check

        Returns:
            File type category (pdf, text, docx, pptx, image)
        """
        ext = Path(filename).suffix.lower()
        if ext == ".pdf":
            return "pdf"
        elif ext == ".txt":
            return "text"
        elif ext in [".docx", ".doc"]:
            return "docx"
        elif ext in [".pptx", ".ppt"]:
            return "pptx"
        elif ext in [".png", ".jpg", ".jpeg"]:
            return "image"
        else:
            return "unknown"

    def process_file(
        self,
        content: bytes,
        filename: str,
    ) -> dict[str, Any]:
        """Process a file and return chunks with metadata.

        Args:
            content: Raw file content as bytes
            filename: Original filename with extension

        Returns:
            Dict containing filename, file_type, chunks, chunk_count,
            detected_language, and total_characters

        Raises:
            ValueError: If file type is not supported
        """
        ext = Path(filename).suffix.lower()

        if not self.is_supported(filename):
            msg = f"Unsupported file type: {ext}. Supported: {self.SUPPORTED_EXTENSIONS}"
            raise ValueError(msg)

        if ext == ".pdf":
            return self._process_pdf(content, filename)
        elif ext == ".txt":
            return self._process_text(content, filename)
        elif ext in [".docx", ".doc"]:
            return self._process_docx(content, filename)
        elif ext in [".pptx", ".ppt"]:
            return self._process_pptx(content, filename)
        elif ext in [".png", ".jpg", ".jpeg"]:
            return self._process_image(content, filename)
        else:
            msg = f"Unsupported file type: {ext}"
            raise ValueError(msg)

    def _process_pdf(self, content: bytes, filename: str) -> dict[str, Any]:
        """Process PDF file.

        Args:
            content: PDF file content
            filename: Original filename

        Returns:
            Processed document data with chunks
        """
        from langchain_community.document_loaders import PyPDFLoader

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        try:
            loader = PyPDFLoader(tmp_path)
            documents = loader.load()

            # Add source metadata
            for doc in documents:
                doc.metadata["source_file"] = filename

            chunks = self.text_splitter.split_documents(documents)
            return self._build_result(chunks, filename, "pdf")
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    def _process_text(self, content: bytes, filename: str) -> dict[str, Any]:
        """Process plain text file.

        Args:
            content: Text file content
            filename: Original filename

        Returns:
            Processed document data with chunks
        """
        text = content.decode("utf-8", errors="replace")
        documents = [Document(
            page_content=text,
            metadata={"source_file": filename}
        )]
        chunks = self.text_splitter.split_documents(documents)
        return self._build_result(chunks, filename, "text")

    def _process_docx(self, content: bytes, filename: str) -> dict[str, Any]:
        """Process Word document.

        Args:
            content: DOCX file content
            filename: Original filename

        Returns:
            Processed document data with chunks
        """
        try:
            from docx import Document as DocxDocument
        except ImportError:
            msg = "python-docx not installed. Install with: pip install python-docx"
            raise ImportError(msg)

        doc = DocxDocument(io.BytesIO(content))
        text_content = "\n\n".join([para.text for para in doc.paragraphs if para.text])

        if not text_content.strip():
            msg = "No text content found in Word document"
            raise ValueError(msg)

        documents = [Document(
            page_content=text_content,
            metadata={"source_file": filename}
        )]
        chunks = self.text_splitter.split_documents(documents)
        return self._build_result(chunks, filename, "docx")

    def _process_pptx(self, content: bytes, filename: str) -> dict[str, Any]:
        """Process PowerPoint file.

        Args:
            content: PPTX file content
            filename: Original filename

        Returns:
            Processed document data with chunks
        """
        try:
            from pptx import Presentation
        except ImportError:
            msg = "python-pptx not installed. Install with: pip install python-pptx"
            raise ImportError(msg)

        prs = Presentation(io.BytesIO(content))
        text_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += shape.text + "\n"
            if slide_text.strip() != f"--- Slide {slide_num} ---":
                text_content.append(slide_text)

        full_text = "\n".join(text_content)

        if not full_text.strip():
            msg = "No text content found in PowerPoint presentation"
            raise ValueError(msg)

        documents = [Document(
            page_content=full_text,
            metadata={"source_file": filename, "total_slides": len(prs.slides)}
        )]
        chunks = self.text_splitter.split_documents(documents)
        return self._build_result(chunks, filename, "pptx")

    def _process_image(self, content: bytes, filename: str) -> dict[str, Any]:
        """Process image with OCR.

        Args:
            content: Image file content
            filename: Original filename

        Returns:
            Processed document data with chunks

        Raises:
            ImportError: If pytesseract or Pillow not installed
            ValueError: If no text extracted from image
        """
        try:
            import pytesseract
            from PIL import Image
        except ImportError as e:
            msg = "pytesseract and Pillow required for image OCR. "
            msg += "Install with: pip install pytesseract Pillow"
            raise ImportError(msg) from e

        try:
            image = Image.open(io.BytesIO(content))
            text = pytesseract.image_to_string(image)
        except Exception as e:
            msg = f"OCR failed: {e}. Ensure Tesseract is installed on your system."
            raise RuntimeError(msg) from e

        if not text.strip():
            msg = "No text could be extracted from the image via OCR"
            raise ValueError(msg)

        documents = [Document(
            page_content=text,
            metadata={
                "source_file": filename,
                "image_size": f"{image.width}x{image.height}",
                "image_mode": image.mode,
            }
        )]
        chunks = self.text_splitter.split_documents(documents)
        return self._build_result(chunks, filename, "image")

    def _build_result(
        self,
        chunks: list[Document],
        filename: str,
        file_type: str,
    ) -> dict[str, Any]:
        """Build standardized result dictionary with language detection.

        Args:
            chunks: List of document chunks
            filename: Original filename
            file_type: Type of file processed

        Returns:
            Standardized result dict with metadata
        """
        # Detect language from first chunk
        sample_text = chunks[0].page_content[:1000] if chunks else ""
        language = self._detect_language(sample_text)

        # Add chunk index to metadata
        for i, chunk in enumerate(chunks):
            chunk.metadata["chunk_index"] = i
            chunk.metadata["file_type"] = file_type

        return {
            "filename": filename,
            "file_type": file_type,
            "chunks": chunks,
            "chunk_count": len(chunks),
            "detected_language": language,
            "total_characters": sum(len(c.page_content) for c in chunks),
        }

    def _detect_language(self, text: str) -> str:
        """Detect language of text.

        Args:
            text: Text to analyze

        Returns:
            Language code (e.g., 'en', 'fr', 'de') or 'unknown'
        """
        if not text.strip():
            return "unknown"

        try:
            from langdetect import detect
            return detect(text)
        except ImportError:
            logger.warning("langdetect not installed, defaulting to 'en'")
            return "en"
        except Exception:
            return "unknown"
